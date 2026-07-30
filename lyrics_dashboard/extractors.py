from __future__ import annotations

import io
import re
from pathlib import Path

from .errors import ExtractionError
from .text_processing import remove_leading_links

TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm"}
SUPPORTED_EXTENSIONS = {
    ".docx",
    ".pdf",
    ".pptx",
    ".xlsx",
    *TEXT_EXTENSIONS,
}


def _decode_text(data: bytes) -> str:
    if data.startswith((b"\xff\xfe", b"\xfe\xff")):
        try:
            return data.decode("utf-16")
        except UnicodeDecodeError:
            pass

    if data and len(data) % 2 == 0:
        even_nuls = data[0::2].count(0)
        odd_nuls = data[1::2].count(0)
        minimum_nuls = max(2, len(data) // 8)
        if odd_nuls >= minimum_nuls and odd_nuls > even_nuls * 4:
            try:
                return data.decode("utf-16-le")
            except UnicodeDecodeError:
                pass
        if even_nuls >= minimum_nuls and even_nuls > odd_nuls * 4:
            try:
                return data.decode("utf-16-be")
            except UnicodeDecodeError:
                pass

    for encoding in ("utf-8-sig", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ExtractionError("The text file encoding could not be detected.")


def _extract_docx(data: bytes) -> tuple[str, frozenset[int]]:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(io.BytesIO(data))
    lines: list[str] = []
    hyperlink_line_indices: set[int] = set()
    hyperlink_tag = qn("w:hyperlink")
    simple_field_tag = qn("w:fldSimple")
    field_instruction_attribute = qn("w:instr")
    instruction_text_tag = qn("w:instrText")

    def is_hyperlink_paragraph(paragraph: Paragraph) -> bool:
        if paragraph.hyperlinks:
            return True

        complex_field_instruction: list[str] = []
        for element in paragraph._p.iter():
            if element.tag == hyperlink_tag:
                return True
            if element.tag == simple_field_tag:
                instruction = element.get(field_instruction_attribute, "")
                if re.search(r"\bHYPERLINK\b", instruction, re.IGNORECASE):
                    return True
            if element.tag == instruction_text_tag:
                complex_field_instruction.append(element.text or "")

        return bool(
            re.search(
                r"\bHYPERLINK\b",
                "".join(complex_field_instruction),
                re.IGNORECASE,
            )
        )

    def append_paragraph(paragraph: Paragraph) -> None:
        paragraph_lines = (
            paragraph.text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
        )
        first_line_index = len(lines)
        lines.extend(paragraph_lines)
        if is_hyperlink_paragraph(paragraph):
            hyperlink_line_indices.update(range(first_line_index, len(lines)))

    for item in document.iter_inner_content():
        if isinstance(item, Paragraph):
            append_paragraph(item)
        elif isinstance(item, Table):
            for row in item.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        append_paragraph(paragraph)
    return "\n".join(lines), frozenset(hyperlink_line_indices)


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception as exc:  # pragma: no cover - library-specific
            raise ExtractionError("Encrypted PDFs are not supported.") from exc
    pages = [(page.extract_text() or "") for page in reader.pages]
    text = "\n".join(pages)
    if not text.strip():
        raise ExtractionError(
            "No selectable text was found in the PDF. Scanned PDFs need OCR, which is not enabled in this MVP."
        )
    return text


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.extend(shape.text.splitlines())
    return "\n".join(lines)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    for worksheet in workbook.worksheets:
        for row in worksheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                lines.append(" ".join(values))
    return "\n".join(lines)


def extract_text(filename: str, data: bytes) -> str:
    """Extract ordered text from a supported upload."""
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ExtractionError(f"Unsupported file type '{suffix or 'unknown'}'. Supported types: {supported}")

    hyperlink_line_indices: frozenset[int] = frozenset()
    try:
        if suffix in TEXT_EXTENSIONS:
            text = _decode_text(data)
        elif suffix == ".docx":
            text, hyperlink_line_indices = _extract_docx(data)
        elif suffix == ".pdf":
            text = _extract_pdf(data)
        elif suffix == ".pptx":
            text = _extract_pptx(data)
        elif suffix == ".xlsx":
            text = _extract_xlsx(data)
        else:  # Defensive; SUPPORTED_EXTENSIONS is exhaustive.
            raise ExtractionError(f"No extractor is configured for '{suffix}'.")
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Could not read '{filename}': {exc}") from exc

    text = remove_leading_links(text, hyperlink_line_indices)
    if not text.strip():
        raise ExtractionError("The uploaded file did not contain readable text.")
    return text
